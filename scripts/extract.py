#!/usr/bin/env python3
"""
Universal document text extractor for translation review.
Supports: xlsx, docx, pptx, pdf, md, html, txt

Usage:
    python extract.py <file> [--source]   # --source flag for Chinese source (skips CJK check)

Output: JSON with sections (text + location) and automated CJK/sensitive-word findings.
"""

import json
import re
import sys
from pathlib import Path


# ─── CJK and sensitive-word detection ───

def has_chinese(text):
    if not text or not isinstance(text, str):
        return False
    return bool(re.search(r'[一-鿿]', text))


SENSITIVE_PATTERNS = [
    (r'\bTaiwan\b(?!\s+Strait)', 'P0', 'Possible Taiwan reference'),
    (r'\bTibet\b', 'P0', 'Possible Tibet reference'),
    (r'\b(drug|narcotic|cocaine|heroin)\b', 'P1', 'Sensitive term'),
    (r'\b(kill|destroy|attack|bomb)\b', 'P1', 'Violence-related term'),
    (r'\b(fraud|scam|illegal)\b', 'P1', 'Legal sensitivity'),
]


def check_sensitive(text):
    flagged = []
    for pattern, severity, desc in SENSITIVE_PATTERNS:
        for m in re.finditer(pattern, text, re.IGNORECASE):
            flagged.append({"term": m.group(), "severity": severity, "description": desc})
    return flagged


# ─── Format parsers ───

def parse_xlsx(path, is_source):
    import openpyxl
    wb = openpyxl.load_workbook(path, data_only=True)
    sections = []
    for sn in wb.sheetnames:
        ws = wb[sn]
        for row in ws.iter_rows(min_row=1, max_row=ws.max_row, max_col=ws.max_column):
            for cell in row:
                if cell.value and str(cell.value).strip():
                    sections.append({
                        "location": f"{sn}!{cell.coordinate}",
                        "text": str(cell.value).strip(),
                    })
    wb.close()
    return sections


def parse_docx(path, is_source):
    import docx
    doc = docx.Document(path)
    sections = []

    # Paragraphs
    for i, para in enumerate(doc.paragraphs):
        text = para.text.strip()
        if text:
            style = para.style.name if para.style else "Normal"
            sections.append({
                "location": f"Paragraph {i+1}" + (f" [{style}]" if style != "Normal" else ""),
                "text": text,
            })

    # Tables
    for ti, table in enumerate(doc.tables):
        for ri, row in enumerate(table.rows):
            for ci, cell in enumerate(row.cells):
                text = cell.text.strip()
                if text:
                    sections.append({
                        "location": f"Table {ti+1}, Row {ri+1}, Col {ci+1}",
                        "text": text,
                    })

    return sections


def parse_pptx(path, is_source):
    import pptx
    prs = pptx.Presentation(path)
    sections = []

    for si, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for pi, para in enumerate(shape.text_frame.paragraphs):
                    text = para.text.strip()
                    if text:
                        sections.append({
                            "location": f"Slide {si+1}, Paragraph {pi+1}",
                            "text": text,
                        })
            if shape.has_table:
                table = shape.table
                for ri, row in enumerate(table.rows):
                    for ci, cell in enumerate(row.cells):
                        text = cell.text.strip()
                        if text:
                            sections.append({
                                "location": f"Slide {si+1}, Table Row {ri+1}, Col {ci+1}",
                                "text": text,
                            })

    return sections


def parse_pdf(path, is_source):
    import pdfplumber
    sections = []
    with pdfplumber.open(path) as pdf:
        for pi, page in enumerate(pdf.pages):
            text = page.extract_text()
            if text:
                for li, line in enumerate(text.split('\n')):
                    line = line.strip()
                    if line and len(line) > 1:
                        sections.append({
                            "location": f"Page {pi+1}, Line {li+1}",
                            "text": line,
                        })
            # Tables
            for table in page.extract_tables():
                for ri, row in enumerate(table):
                    for ci, cell in enumerate(row):
                        if cell and cell.strip():
                            sections.append({
                                "location": f"Page {pi+1}, Table Row {ri+1}, Col {ci+1}",
                                "text": cell.strip(),
                            })
    return sections


def parse_md(path, is_source):
    sections = []
    with open(path, 'r', encoding='utf-8') as f:
        lines = f.readlines()
    in_code = False
    for li, line in enumerate(lines):
        stripped = line.strip()
        if stripped.startswith('```'):
            in_code = not in_code
            continue
        if in_code:
            continue
        if stripped and not stripped.startswith('#') and len(stripped) > 1:
            sections.append({
                "location": f"Line {li+1}",
                "text": stripped.lstrip('#').strip(),
            })
    return sections


def parse_html(path, is_source):
    from bs4 import BeautifulSoup
    with open(path, 'r', encoding='utf-8') as f:
        soup = BeautifulSoup(f.read(), 'lxml')

    # Remove script/style
    for tag in soup(['script', 'style', 'nav', 'footer', 'code', 'pre']):
        tag.decompose()

    sections = []
    # Headings
    for tag in soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'li', 'td', 'th']):
        text = tag.get_text(strip=True)
        if text and len(text) > 1:
            sections.append({
                "location": f"<{tag.name}> \"{text[:40]}...\"" if len(text) > 40 else f"<{tag.name}> \"{text}\"",
                "text": text,
            })

    return sections


def parse_txt(path, is_source):
    sections = []
    with open(path, 'r', encoding='utf-8') as f:
        for li, line in enumerate(f):
            stripped = line.strip()
            if stripped and len(stripped) > 1:
                sections.append({
                    "location": f"Line {li+1}",
                    "text": stripped,
                })
    return sections


# ─── Format dispatch ───

PARSERS = {
    '.xlsx': parse_xlsx,
    '.xls': parse_xlsx,
    '.docx': parse_docx,
    '.doc': parse_docx,
    '.pptx': parse_pptx,
    '.ppt': parse_pptx,
    '.pdf': parse_pdf,
    '.md': parse_md,
    '.markdown': parse_md,
    '.html': parse_html,
    '.htm': parse_html,
    '.txt': parse_txt,
    '.csv': parse_txt,
}


def main():
    if len(sys.argv) < 2:
        print("Usage: python extract.py <file> [--source]", file=sys.stderr)
        sys.exit(1)

    path = Path(sys.argv[1])
    if not path.exists():
        print(f"Error: file not found: {path}", file=sys.stderr)
        sys.exit(1)

    is_source = '--source' in sys.argv
    ext = path.suffix.lower()
    parser = PARSERS.get(ext, parse_txt)

    sections = parser(str(path), is_source)

    # Automated checks (only for target files, not source)
    findings = []
    if not is_source:
        for sec in sections:
            if has_chinese(sec["text"]):
                findings.append({
                    "dimension": "漏翻",
                    "severity": "P1",
                    "location": sec["location"],
                    "description": f"Untranslated Chinese: '{sec['text'][:60]}'",
                })
            for s in check_sensitive(sec["text"]):
                findings.append({
                    "dimension": "敏感词",
                    "severity": s["severity"],
                    "location": sec["location"],
                    "description": f"{s['description']}: '{s['term']}'",
                })

    # Deduplicate texts for review efficiency
    unique_texts = {}
    for sec in sections:
        t = sec["text"]
        if t not in unique_texts:
            unique_texts[t] = {"text": t, "locations": []}
        unique_texts[t]["locations"].append(sec["location"])

    result = {
        "file": str(path.absolute()),
        "format": ext.lstrip('.'),
        "is_source": is_source,
        "total_sections": len(sections),
        "unique_texts": len(unique_texts),
        "sections": sections,
        "automated_findings": findings,
        "automated_count": len(findings),
    }

    # For target files, provide deduplicated texts for manual review
    if not is_source:
        result["review_texts"] = [
            {"text": v["text"], "locations": v["locations"][:3], "occurrence_count": len(v["locations"])}
            for v in sorted(unique_texts.values(), key=lambda x: -len(x["locations"]))
        ]

    print(json.dumps(result, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
